#!/usr/bin/env python3
"""Build the cleaned idtpptx V1 template PPTX.

The generated deck contains neutral sample slides for registered layout IDs.
It intentionally starts from the existing V1 template so the master chrome,
logo, title band, and theme relationships stay template-derived while the
visible slides are rebuilt from the registered clean layouts.
"""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_TEMPLATE = ROOT / "assets/templates/inspur-pragmatic-template-v1.pptx"
DEFAULT_OUTPUT = DEFAULT_TEMPLATE

FONT = "Microsoft YaHei"

COLOR_PRIMARY_BLUE = "0062AC"
COLOR_DEEP_BLUE = "00518E"
COLOR_DARK_NAVY = "213261"
COLOR_BODY = "111111"
COLOR_LIGHT_FILL = "F2F4F7"
COLOR_MID_FILL = "EDEDED"
COLOR_SOFT_WHITE = "FAFAFA"
COLOR_GRID = "A4A3A4"
COLOR_RED = "D93025"
COLOR_WHITE = "FFFFFF"


def inch(value: float):
    return Inches(value)


def rgb(hex_value: str) -> RGBColor:
    value = hex_value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def set_text(
    shape,
    text: str,
    *,
    font_size: float,
    color: str = COLOR_BODY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin_x: float = 0,
    margin_y: float = 0,
) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = inch(margin_x)
    text_frame.margin_right = inch(margin_x)
    text_frame.margin_top = inch(margin_y)
    text_frame.margin_bottom = inch(margin_y)
    text_frame.vertical_anchor = valign

    lines = text.split("\n")
    for idx, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def style_title(slide, title: str) -> None:
    title_shape = slide.shapes.title
    title_shape.text = ""
    title_shape.left = inch(0.37)
    title_shape.top = inch(0.10)
    title_shape.width = inch(12.60)
    title_shape.height = inch(0.72)
    set_text(
        title_shape,
        title,
        font_size=30,
        color=COLOR_DARK_NAVY,
        bold=True,
        align=PP_ALIGN.LEFT,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_textbox(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: float,
    color: str = COLOR_BODY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin_x: float = 0,
    margin_y: float = 0,
):
    shape = slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
    set_text(
        shape,
        text,
        font_size=font_size,
        color=color,
        bold=bold,
        align=align,
        valign=valign,
        margin_x=margin_x,
        margin_y=margin_y,
    )
    return shape


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = COLOR_SOFT_WHITE,
    line: str = COLOR_GRID,
    line_width: float = 0.75,
    text: str = "",
    font_size: float = 12,
    text_color: str = COLOR_BODY,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
    valign=MSO_ANCHOR.MIDDLE,
    margin_x: float = 0,
    margin_y: float = 0,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(line_width)
    if text:
        set_text(
            shape,
            text,
            font_size=font_size,
            color=text_color,
            bold=bold,
            align=align,
            valign=valign,
            margin_x=margin_x,
            margin_y=margin_y,
        )
    return shape


def add_arrow(slide, x: float, y: float, w: float, h: float, *, color: str = COLOR_RED):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, inch(x), inch(y), inch(w), inch(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)
    return shape


def add_table_cell(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str,
    line: str = COLOR_GRID,
    line_width: float = 0.45,
    font_size: float = 12,
    text_color: str = COLOR_BODY,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    return add_rect(
        slide,
        x,
        y,
        w,
        h,
        fill=fill,
        line=line,
        line_width=line_width,
        text=text,
        font_size=font_size,
        text_color=text_color,
        bold=bold,
        align=align,
        valign=MSO_ANCHOR.MIDDLE,
        margin_x=0.08,
        margin_y=0.04,
    )


def remove_existing_slides(prs: Presentation) -> None:
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
    for slide_id in list(slide_id_list):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def add_cov_02a(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    for placeholder in list(slide.placeholders):
        placeholder._element.getparent().remove(placeholder._element)  # noqa: SLF001 - remove inherited cover prompts.

    # Cover v1 keeps the inherited blue top header, but clears the old gray
    # middle band so the whole body reads as a white canvas.
    white_body = add_rect(
        slide,
        0.0,
        1.74,
        13.333,
        4.18,
        fill=COLOR_WHITE,
        line=COLOR_WHITE,
        line_width=0,
    )
    white_body.name = "Cover White Body Background"

    add_textbox(
        slide,
        "技术方案评审材料",
        2.30,
        2.88,
        8.70,
        0.70,
        font_size=30,
        color=COLOR_DARK_NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    add_textbox(
        slide,
        "背景、结论与后续安排",
        2.55,
        3.62,
        8.20,
        0.36,
        font_size=15,
        color=COLOR_BODY,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        "汇报部门：业务技术团队",
        0.95,
        4.95,
        4.40,
        0.34,
        font_size=13,
        color=COLOR_BODY,
    )
    add_textbox(
        slide,
        "日期：2026 年 5 月",
        0.95,
        5.32,
        4.40,
        0.34,
        font_size=13,
        color=COLOR_BODY,
    )


def add_dir_01a(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_title(slide, "目录进度示例")

    add_textbox(
        slide,
        "长培训或方案说明先给读者定位当前位置。",
        0.70,
        1.12,
        11.95,
        0.36,
        font_size=13,
        color=COLOR_BODY,
    )

    rows = [
        ("1", "背景和目标", "先说明为什么要做"),
        ("2", "方案设计", "本页高亮当前章节"),
        ("3", "实施路径", "再说明怎么推进"),
        ("4", "验证结果", "最后沉淀结论和动作"),
    ]
    x0 = 2.30
    y0 = 1.82
    row_w = 8.75
    row_h = 0.56
    gap = 0.24
    for idx, (number, title, note) in enumerate(rows):
        y = y0 + idx * (row_h + gap)
        active = idx == 1
        fill = COLOR_PRIMARY_BLUE if active else COLOR_LIGHT_FILL
        text_color = COLOR_WHITE if active else COLOR_DARK_NAVY
        line = COLOR_PRIMARY_BLUE if active else COLOR_GRID

        add_rect(
            slide,
            x0,
            y,
            0.56,
            row_h,
            fill=fill,
            line=line,
            line_width=0.6,
            text=number,
            font_size=14,
            text_color=text_color,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_rect(
            slide,
            x0 + 0.70,
            y,
            row_w - 0.70,
            row_h,
            fill=fill if active else COLOR_SOFT_WHITE,
            line=line,
            line_width=0.6,
        )
        add_textbox(
            slide,
            title,
            x0 + 0.95,
            y + 0.10,
            2.20,
            0.32,
            font_size=13,
            color=text_color,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            note,
            x0 + 3.35,
            y + 0.10,
            5.30,
            0.32,
            font_size=11.5,
            color=text_color if active else COLOR_BODY,
            valign=MSO_ANCHOR.MIDDLE,
        )

    add_rect(
        slide,
        0.70,
        5.95,
        11.95,
        0.58,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.6,
        text="使用：超过 12 页或跨多个章节时，用目录页承接阅读节奏。",
        font_size=11.5,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.14,
        margin_y=0.06,
    )


def add_prc_03a(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_title(slide, "阶段推进示例")

    add_textbox(
        slide,
        "先试点验证，再按范围逐步推广。",
        0.70,
        1.12,
        11.95,
        0.36,
        font_size=13,
        color=COLOR_BODY,
    )

    axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, inch(1.05), inch(3.00), inch(11.10), inch(0.06))
    axis.fill.solid()
    axis.fill.fore_color.rgb = rgb(COLOR_PRIMARY_BLUE)
    axis.line.color.rgb = rgb(COLOR_PRIMARY_BLUE)
    axis.line.width = Pt(0)

    phases = [
        ("阶段一", "调研确认", "明确对象、口径和边界"),
        ("阶段二", "试点验证", "选择代表场景做闭环"),
        ("阶段三", "推广落地", "扩展到主要团队和项目"),
        ("阶段四", "持续优化", "按数据反馈迭代规则"),
    ]
    x_positions = [0.90, 3.80, 6.70, 9.60]

    for x, (phase, title, body) in zip(x_positions, phases):
        add_rect(
            slide,
            x,
            1.85,
            2.55,
            2.70,
            fill=COLOR_SOFT_WHITE,
            line=COLOR_GRID,
            line_width=0.75,
        )
        add_rect(
            slide,
            x,
            1.85,
            2.55,
            0.42,
            fill=COLOR_PRIMARY_BLUE,
            line=COLOR_PRIMARY_BLUE,
            text=phase,
            font_size=11,
            text_color=COLOR_WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x + 1.155), inch(2.88), inch(0.24), inch(0.24))
        marker.fill.solid()
        marker.fill.fore_color.rgb = rgb(COLOR_PRIMARY_BLUE)
        marker.line.color.rgb = rgb(COLOR_PRIMARY_BLUE)
        add_textbox(
            slide,
            title,
            x + 0.18,
            3.28,
            2.19,
            0.34,
            font_size=14,
            color=COLOR_DARK_NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            body,
            x + 0.20,
            3.78,
            2.15,
            0.54,
            font_size=11.5,
            color=COLOR_BODY,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )

    add_rect(
        slide,
        0.70,
        5.72,
        11.95,
        0.62,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.6,
        text="依赖：每阶段结束后先确认口径，再进入下一阶段。",
        font_size=11.5,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.14,
        margin_y=0.06,
    )


def add_ss_02a(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_title(slide, "操作步骤示例")

    add_textbox(
        slide,
        "按入口、动作、结果组织截图，重点控件用红框标注。",
        0.70,
        1.12,
        11.95,
        0.36,
        font_size=13,
        color=COLOR_BODY,
    )

    for x, label in ((0.70, "界面示意一"), (7.20, "界面示意二")):
        add_rect(
            slide,
            x,
            1.72,
            5.45,
            3.40,
            fill=COLOR_LIGHT_FILL,
            line=COLOR_GRID,
            line_width=0.9,
        )
        add_rect(
            slide,
            x + 0.30,
            2.22,
            4.85,
            0.42,
            fill=COLOR_MID_FILL,
            line=COLOR_GRID,
            line_width=0.4,
            text=label,
            font_size=10.5,
            text_color=COLOR_DARK_NAVY,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_rect(
            slide,
            x + 0.30,
            3.00,
            4.85,
            1.32,
            fill=COLOR_MID_FILL,
            line=COLOR_GRID,
            line_width=0.4,
        )
        add_textbox(
            slide,
            "真实截图放入此区域，保持同一比例",
            x + 0.52,
            3.46,
            2.85,
            0.36,
            font_size=10.5,
            color=COLOR_BODY,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_rect(
            slide,
            x + 3.55,
            3.12,
            1.20,
            0.62,
            fill=COLOR_LIGHT_FILL,
            line=COLOR_RED,
            line_width=1.2,
        )

    add_arrow(slide, 6.34, 3.10, 0.60, 0.30)

    add_textbox(
        slide,
        "第一步：进入页面",
        0.70,
        5.22,
        5.45,
        0.36,
        font_size=11,
        color=COLOR_BODY,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        "第二步：确认结果",
        7.20,
        5.22,
        5.45,
        0.36,
        font_size=11,
        color=COLOR_BODY,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_rect(
        slide,
        0.70,
        5.95,
        11.95,
        0.58,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.6,
        text="说明：真实截图交付前需脱敏，截图文字必须可读。",
        font_size=11.5,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.14,
        margin_y=0.06,
    )


def add_ss_03a(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_title(slide, "故障排查示例")

    add_textbox(
        slide,
        "先定位异常提示，再说明判断依据和处理动作。",
        0.70,
        1.12,
        11.95,
        0.36,
        font_size=13,
        color=COLOR_BODY,
    )

    add_rect(
        slide,
        0.70,
        1.72,
        7.25,
        3.95,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.9,
    )
    add_rect(
        slide,
        1.05,
        2.10,
        6.55,
        0.42,
        fill=COLOR_MID_FILL,
        line=COLOR_GRID,
        line_width=0.4,
        text="异常提示区域",
        font_size=11,
        text_color=COLOR_DARK_NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_rect(
        slide,
        1.05,
        2.88,
        6.55,
        1.62,
        fill=COLOR_MID_FILL,
        line=COLOR_GRID,
        line_width=0.4,
    )
    add_textbox(
        slide,
        "真实截图放入此区域，保留异常提示和当前页面。",
        1.35,
        3.44,
        4.15,
        0.32,
        font_size=11,
        color=COLOR_BODY,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_rect(
        slide,
        5.80,
        3.04,
        1.35,
        0.68,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_RED,
        line_width=1.2,
    )
    add_textbox(
        slide,
        "截图证据",
        0.70,
        5.22,
        7.25,
        0.34,
        font_size=11,
        color=COLOR_BODY,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    diagnosis_blocks = [
        (1.72, "现象", "页面提示保存失败"),
        (3.12, "判断", "先确认配置和权限"),
        (4.52, "处理", "修正后重新提交"),
    ]
    for y, heading, body in diagnosis_blocks:
        add_rect(
            slide,
            8.35,
            y,
            4.30,
            1.15,
            fill=COLOR_SOFT_WHITE,
            line=COLOR_GRID,
            line_width=0.75,
        )
        add_rect(
            slide,
            8.35,
            y,
            0.12,
            1.15,
            fill=COLOR_PRIMARY_BLUE,
            line=COLOR_PRIMARY_BLUE,
            line_width=0,
        )
        add_textbox(
            slide,
            heading,
            8.62,
            y + 0.16,
            3.75,
            0.28,
            font_size=13,
            color=COLOR_DARK_NAVY,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            body,
            8.62,
            y + 0.58,
            3.75,
            0.38,
            font_size=11.5,
            color=COLOR_BODY,
            valign=MSO_ANCHOR.MIDDLE,
        )

    add_rect(
        slide,
        0.70,
        5.95,
        11.95,
        0.58,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.6,
        text="边界：本页只说明单个异常，复杂问题拆成多页排查。",
        font_size=11.5,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.14,
        margin_y=0.06,
    )


def add_arc_01a(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_title(slide, "架构说明示例")

    add_textbox(
        slide,
        "左侧先给阅读口径，右侧展示主链路。",
        0.70,
        1.12,
        11.95,
        0.36,
        font_size=13,
        color=COLOR_BODY,
    )

    add_rect(slide, 0.70, 1.65, 3.45, 4.42, fill=COLOR_SOFT_WHITE, line=COLOR_GRID, line_width=0.75)
    add_rect(slide, 0.70, 1.65, 0.12, 4.42, fill=COLOR_PRIMARY_BLUE, line=COLOR_PRIMARY_BLUE, line_width=0)
    add_textbox(
        slide,
        "阅读口径",
        1.02,
        1.92,
        2.80,
        0.34,
        font_size=14,
        color=COLOR_DARK_NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide,
        "明确对象边界\n确认主链路方向\n标出异常入口\n写清交付结论",
        1.02,
        2.50,
        2.78,
        1.45,
        font_size=12,
        color=COLOR_BODY,
        valign=MSO_ANCHOR.TOP,
    )
    add_rect(
        slide,
        1.02,
        4.55,
        2.78,
        0.78,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.5,
        text="边界：只讲和本页结论相关的链路。",
        font_size=10.8,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.10,
        margin_y=0.06,
    )

    add_rect(slide, 4.55, 1.65, 8.10, 4.42, fill=COLOR_SOFT_WHITE, line=COLOR_GRID, line_width=0.75)
    add_textbox(
        slide,
        "主链路示意",
        4.90,
        1.92,
        7.35,
        0.30,
        font_size=13,
        color=COLOR_DARK_NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    blocks = [
        (5.00, 2.70, "用户入口", COLOR_LIGHT_FILL, COLOR_DARK_NAVY),
        (6.90, 2.70, "服务网关", COLOR_PRIMARY_BLUE, COLOR_WHITE),
        (8.80, 2.70, "核心服务", COLOR_DEEP_BLUE, COLOR_WHITE),
        (10.70, 2.70, "数据存储", COLOR_LIGHT_FILL, COLOR_DARK_NAVY),
    ]
    for x, y, label, fill, text_color in blocks:
        add_rect(
            slide,
            x,
            y,
            1.34,
            0.62,
            fill=fill,
            line=COLOR_GRID if fill == COLOR_LIGHT_FILL else fill,
            text=label,
            font_size=11.5,
            text_color=text_color,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    for x in (6.42, 8.32, 10.22):
        add_arrow(slide, x, 2.90, 0.34, 0.18, color=COLOR_PRIMARY_BLUE)

    add_rect(
        slide,
        6.40,
        4.12,
        1.58,
        0.54,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        text="任务队列",
        font_size=10.8,
        text_color=COLOR_DARK_NAVY,
        bold=True,
    )
    add_rect(
        slide,
        8.60,
        4.12,
        1.58,
        0.54,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        text="审计日志",
        font_size=10.8,
        text_color=COLOR_DARK_NAVY,
        bold=True,
    )
    add_rect(
        slide,
        5.00,
        5.10,
        7.30,
        0.42,
        fill=COLOR_MID_FILL,
        line=COLOR_GRID,
        text="说明：复杂链路先拆主干，异常分支放到下一页。",
        font_size=10.8,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.10,
        margin_y=0.04,
    )

    add_rect(
        slide,
        0.70,
        6.24,
        11.95,
        0.36,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.5,
        text="适用：说明文字和架构图同等重要；图特别宽时改用 ARC-02。",
        font_size=10.8,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.14,
        margin_y=0.03,
    )


def add_arc_03a(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_title(slide, "图文拆分架构示例")

    add_textbox(
        slide,
        "先看系统关系，再读右侧约束。",
        0.70,
        1.12,
        11.95,
        0.36,
        font_size=13,
        color=COLOR_BODY,
    )

    add_rect(slide, 0.70, 1.62, 7.15, 4.55, fill=COLOR_SOFT_WHITE, line=COLOR_GRID, line_width=0.75)
    add_textbox(
        slide,
        "架构关系",
        1.02,
        1.88,
        6.45,
        0.30,
        font_size=13,
        color=COLOR_DARK_NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    top_nodes = [
        (1.05, "入口层", COLOR_LIGHT_FILL, COLOR_DARK_NAVY),
        (3.00, "编排服务", COLOR_PRIMARY_BLUE, COLOR_WHITE),
        (4.95, "能力组件", COLOR_DEEP_BLUE, COLOR_WHITE),
        (6.50, "数据底座", COLOR_LIGHT_FILL, COLOR_DARK_NAVY),
    ]
    for x, label, fill, text_color in top_nodes:
        add_rect(
            slide,
            x,
            2.42,
            1.15,
            0.58,
            fill=fill,
            line=COLOR_GRID if fill == COLOR_LIGHT_FILL else fill,
            text=label,
            font_size=10.8,
            text_color=text_color,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    for x in (2.33, 4.28, 6.18):
        add_arrow(slide, x, 2.61, 0.34, 0.18, color=COLOR_PRIMARY_BLUE)

    add_rect(slide, 1.05, 3.50, 6.45, 1.38, fill=COLOR_LIGHT_FILL, line=COLOR_GRID, line_width=0.6)
    add_textbox(
        slide,
        "能力分层",
        1.28,
        3.68,
        1.30,
        0.26,
        font_size=10.8,
        color=COLOR_DARK_NAVY,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    layer_items = [
        (2.15, "任务调度"),
        (3.72, "状态同步"),
        (5.29, "审计追踪"),
    ]
    for x, label in layer_items:
        add_rect(
            slide,
            x,
            3.98,
            1.25,
            0.42,
            fill=COLOR_SOFT_WHITE,
            line=COLOR_GRID,
            line_width=0.5,
            text=label,
            font_size=9.8,
            text_color=COLOR_BODY,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )

    add_rect(
        slide,
        1.05,
        5.24,
        6.45,
        0.46,
        fill=COLOR_MID_FILL,
        line=COLOR_GRID,
        line_width=0.5,
        text="示意图只保留主干结构，细节放到后续页展开。",
        font_size=10.5,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.10,
        margin_y=0.04,
    )

    add_rect(slide, 8.15, 1.62, 4.50, 4.55, fill=COLOR_SOFT_WHITE, line=COLOR_GRID, line_width=0.75)
    add_textbox(
        slide,
        "说明口径",
        8.45,
        1.88,
        3.90,
        0.30,
        font_size=13,
        color=COLOR_DARK_NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )

    cards = [
        ("入口边界", "统一入口、权限和对象范围"),
        ("核心链路", "先编排任务，再调用能力组件"),
        ("数据约束", "关键状态和结果需要可追溯"),
        ("输出结果", "面向评审沉淀结论和行动"),
    ]
    y = 2.42
    for title, body in cards:
        add_rect(slide, 8.45, y, 3.90, 0.64, fill=COLOR_LIGHT_FILL, line=COLOR_GRID, line_width=0.5)
        add_rect(slide, 8.45, y, 0.08, 0.64, fill=COLOR_PRIMARY_BLUE, line=COLOR_PRIMARY_BLUE, line_width=0)
        add_textbox(
            slide,
            title,
            8.64,
            y + 0.08,
            3.48,
            0.20,
            font_size=10.8,
            color=COLOR_DARK_NAVY,
            bold=True,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            body,
            8.64,
            y + 0.34,
            3.48,
            0.18,
            font_size=9.6,
            color=COLOR_BODY,
            valign=MSO_ANCHOR.MIDDLE,
        )
        y += 0.78

    add_rect(
        slide,
        0.70,
        6.28,
        10.65,
        0.34,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.5,
        text="适用：图和文字同等重要；可按阅读顺序左右互换。",
        font_size=10.5,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.14,
        margin_y=0.03,
    )


def add_arc_02a(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_title(slide, "技术架构示例")

    add_textbox(
        slide,
        "先看主链路，再看边界条件。",
        0.70,
        1.12,
        11.95,
        0.36,
        font_size=13,
        color=COLOR_BODY,
    )

    add_rect(slide, 0.70, 1.62, 11.95, 4.25, fill=COLOR_SOFT_WHITE, line=COLOR_GRID, line_width=0.8)

    add_rect(
        slide,
        1.05,
        2.08,
        2.00,
        0.82,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        text="外部系统",
        font_size=13,
        text_color=COLOR_DARK_NAVY,
        bold=True,
    )
    add_rect(
        slide,
        3.65,
        2.08,
        2.00,
        0.82,
        fill=COLOR_PRIMARY_BLUE,
        line=COLOR_PRIMARY_BLUE,
        text="入口层",
        font_size=13,
        text_color=COLOR_WHITE,
        bold=True,
    )
    add_rect(
        slide,
        6.25,
        2.08,
        2.00,
        0.82,
        fill=COLOR_DEEP_BLUE,
        line=COLOR_DEEP_BLUE,
        text="服务层",
        font_size=13,
        text_color=COLOR_WHITE,
        bold=True,
    )
    add_rect(
        slide,
        8.85,
        2.08,
        2.00,
        0.82,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        text="数据层",
        font_size=13,
        text_color=COLOR_DARK_NAVY,
        bold=True,
    )

    for x in (3.15, 5.75, 8.35):
        add_arrow(slide, x, 2.32, 0.38, 0.22, color=COLOR_PRIMARY_BLUE)

    add_rect(
        slide,
        1.05,
        3.52,
        8.10,
        1.08,
        fill=COLOR_MID_FILL,
        line=COLOR_GRID,
        text="主链路说明区：用简短中文标签说明数据、任务或请求如何流转。",
        font_size=12,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
    )
    add_rect(
        slide,
        9.70,
        3.72,
        0.92,
        0.48,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_RED,
        line_width=1.2,
        text="边界",
        font_size=10.5,
        text_color=COLOR_RED,
        bold=True,
    )

    add_rect(
        slide,
        0.70,
        6.05,
        11.95,
        0.55,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.6,
        text="边界：本页只展示主要链路，不展开异常分支。",
        font_size=11.5,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.14,
        margin_y=0.06,
    )


def add_tbl_02a(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_title(slide, "证据表说明示例")

    add_textbox(
        slide,
        "密集表格先统一口径，再压缩到可读范围。",
        0.70,
        1.12,
        11.95,
        0.36,
        font_size=13,
        color=COLOR_BODY,
    )

    x0 = 0.70
    y0 = 1.72
    col_widths = [1.05, 1.50, 2.05, 2.55, 2.20, 2.60]
    row_heights = [0.48, 0.58, 0.58, 0.58, 0.58, 0.58]
    headers = ["类别", "对象", "口径", "证据", "约束", "结论"]
    rows = [
        ["范围", "项目", "同一对象", "来源可追溯", "不跨系统合并", "先确认边界"],
        ["质量", "结果", "同一时间窗", "状态可复核", "异常单独说明", "支撑评审"],
        ["时效", "任务", "按日更新", "处理链路完整", "延迟需标注", "便于跟踪"],
        ["风险", "异常", "分级记录", "处置动作明确", "不可只写现象", "降低误判"],
        ["输出", "结论", "口径一致", "责任人明确", "不填伪数据", "可进入复盘"],
    ]

    y = y0
    x = x0
    for header, width in zip(headers, col_widths):
        add_table_cell(
            slide,
            header,
            x,
            y,
            width,
            row_heights[0],
            fill=COLOR_PRIMARY_BLUE,
            line=COLOR_WHITE,
            line_width=0.6,
            font_size=11.5,
            text_color=COLOR_WHITE,
            bold=True,
            align=PP_ALIGN.LEFT,
        )
        x += width

    y += row_heights[0]
    for row_index, row in enumerate(rows):
        x = x0
        fill = COLOR_WHITE if row_index % 2 == 0 else COLOR_LIGHT_FILL
        for value, width in zip(row, col_widths):
            add_table_cell(
                slide,
                value,
                x,
                y,
                width,
                row_heights[row_index + 1],
                fill=fill,
                line=COLOR_GRID,
                line_width=0.42,
                font_size=10.5,
                text_color=COLOR_BODY,
                align=PP_ALIGN.LEFT,
            )
            x += width
        y += row_heights[row_index + 1]

    add_rect(
        slide,
        0.70,
        5.55,
        11.95,
        0.72,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.6,
        text="规则：超过 6 列或 8 行时优先拆表，不把正文缩到不可读。",
        font_size=11.5,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.14,
        margin_y=0.06,
    )


def add_tbl_03a(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_title(slide, "指标评估示例")

    add_textbox(
        slide,
        "先明确判断口径，再说明通过标准。",
        0.70,
        1.12,
        11.95,
        0.36,
        font_size=13,
        color=COLOR_BODY,
    )

    x0 = 0.70
    y0 = 1.72
    col_widths = [1.60, 3.10, 3.55, 3.70]
    row_heights = [0.52, 0.78, 0.78, 0.78, 0.78]
    headers = ["指标", "判断内容", "通过标准", "业务含义"]
    rows = [
        ["范围", "是否限定对象边界", "同一对象、同一时间窗", "避免跨范围误判"],
        ["质量", "证据是否可解释", "来源、口径、结果可追溯", "支撑评审复核"],
        ["时效", "处理是否及时", "关键状态按日更新", "便于持续跟踪"],
        ["风险", "是否有异常兜底", "异常场景有处理路径", "降低交付风险"],
    ]

    y = y0
    x = x0
    for header, width in zip(headers, col_widths):
        add_table_cell(
            slide,
            header,
            x,
            y,
            width,
            row_heights[0],
            fill=COLOR_PRIMARY_BLUE,
            line=COLOR_WHITE,
            line_width=0.6,
            font_size=13,
            text_color=COLOR_WHITE,
            bold=True,
            align=PP_ALIGN.LEFT,
        )
        x += width

    y += row_heights[0]
    for row_index, row in enumerate(rows):
        x = x0
        fill = COLOR_WHITE if row_index % 2 == 0 else COLOR_LIGHT_FILL
        for value, width in zip(row, col_widths):
            add_table_cell(
                slide,
                value,
                x,
                y,
                width,
                row_heights[row_index + 1],
                fill=fill,
                line=COLOR_GRID,
                line_width=0.45,
                font_size=12,
                text_color=COLOR_BODY,
                align=PP_ALIGN.LEFT,
            )
            x += width
        y += row_heights[row_index + 1]

    add_rect(
        slide,
        0.70,
        6.08,
        11.95,
        0.52,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.6,
        text="输出：通过标准不清晰时，先补口径，不急于填结论。",
        font_size=11.5,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.14,
        margin_y=0.06,
    )


def add_sum_02a(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    style_title(slide, "决策复盘示例")

    add_rect(
        slide,
        0.70,
        1.18,
        11.95,
        0.68,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.6,
        text="结论：方案可以进入试点，先补齐口径和责任边界。",
        font_size=14,
        text_color=COLOR_DARK_NAVY,
        bold=True,
        align=PP_ALIGN.LEFT,
        margin_x=0.16,
        margin_y=0.06,
    )

    blocks = [
        (0.70, "结论", "同意按试点范围推进"),
        (4.84, "责任", "明确负责人和协同团队"),
        (8.98, "下一步", "完成口径确认和材料补齐"),
    ]
    for x, heading, body in blocks:
        add_rect(
            slide,
            x,
            2.16,
            3.65,
            1.42,
            fill=COLOR_SOFT_WHITE,
            line=COLOR_GRID,
            line_width=0.75,
        )
        add_rect(
            slide,
            x,
            2.16,
            3.65,
            0.38,
            fill=COLOR_PRIMARY_BLUE,
            line=COLOR_PRIMARY_BLUE,
            line_width=0,
            text=heading,
            font_size=12.5,
            text_color=COLOR_WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            body,
            x + 0.18,
            2.76,
            3.29,
            0.52,
            font_size=12,
            color=COLOR_BODY,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )

    x0 = 0.70
    y0 = 4.02
    col_widths = [1.80, 2.15, 2.05, 5.95]
    row_heights = [0.34, 0.37, 0.37, 0.37]
    headers = ["事项", "负责人", "时间", "输出"]
    rows = [
        ["口径确认", "业务团队", "本周内", "统一判断标准"],
        ["试点验证", "项目团队", "下阶段", "形成验证结果"],
        ["复盘汇报", "牵头团队", "评审后", "同步风险和结论"],
    ]

    y = y0
    x = x0
    for header, width in zip(headers, col_widths):
        add_table_cell(
            slide,
            header,
            x,
            y,
            width,
            row_heights[0],
            fill=COLOR_PRIMARY_BLUE,
            line=COLOR_WHITE,
            line_width=0.6,
            font_size=11.5,
            text_color=COLOR_WHITE,
            bold=True,
            align=PP_ALIGN.LEFT,
        )
        x += width

    y += row_heights[0]
    for row_index, row in enumerate(rows):
        x = x0
        fill = COLOR_WHITE if row_index % 2 == 0 else COLOR_LIGHT_FILL
        for value, width in zip(row, col_widths):
            add_table_cell(
                slide,
                value,
                x,
                y,
                width,
                row_heights[row_index + 1],
                fill=fill,
                line=COLOR_GRID,
                line_width=0.45,
                font_size=10.8,
                text_color=COLOR_BODY,
                align=PP_ALIGN.LEFT,
            )
            x += width
        y += row_heights[row_index + 1]

    add_rect(
        slide,
        0.70,
        5.88,
        11.95,
        0.58,
        fill=COLOR_LIGHT_FILL,
        line=COLOR_GRID,
        line_width=0.6,
        text="跟进：重要结论进入项目清单，避免只停留在会议纪要。",
        font_size=11.5,
        text_color=COLOR_BODY,
        align=PP_ALIGN.LEFT,
        margin_x=0.14,
        margin_y=0.06,
    )


def build_template(template: Path, output: Path) -> Path:
    prs = Presentation(str(template))
    remove_existing_slides(prs)
    add_cov_02a(prs)
    add_dir_01a(prs)
    add_prc_03a(prs)
    add_ss_02a(prs)
    add_ss_03a(prs)
    add_arc_01a(prs)
    add_arc_03a(prs)
    add_arc_02a(prs)
    add_tbl_02a(prs)
    add_tbl_03a(prs)
    add_sum_02a(prs)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return output


def main() -> int:
    parser = argparse.ArgumentParser(description="Build the idtpptx cleaned V1 template PPTX.")
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()

    output = build_template(args.template.resolve(), args.output.resolve())
    print(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
